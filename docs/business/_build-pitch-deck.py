# -*- coding: utf-8 -*-
"""عرض «أمانة» التقديمي، دراسة جدوى / Pitch Deck.

يُبنى بهوية الموقع نفسها (بنفسجي #7C3AED + وردي #EC4899) وباتجاه من اليمين لليسار.
كل رقم هنا مشتقّ من docs/business/tech-budget-3y.md والنموذج المالي المصحّح.
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r"D:\amana\amana-platform\docs\business\amana-pitch-deck-21-07-2026.pptx"
LOGO = r"D:\amana\amana-platform\apps\admin\public\amana-logo.jpg"  # نفس شعار الموقع

# ─────────── الهوية البصرية، مطابقة لـ apps/admin landing-client.tsx ───────────
INK        = RGBColor(0x1E, 0x14, 0x2E)
PAPER      = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE     = RGBColor(0x6D, 0x28, 0xD9)   # violet-700، العناوين وترويسات الجداول
PURPLE_DEEP= RGBColor(0x4C, 0x1D, 0x95)   # violet-900، الشرائح الداكنة
PURPLE_MID = RGBColor(0x7C, 0x3A, 0xED)   # violet-600، لون الموقع الأساسي
PURPLE_PALE= RGBColor(0xF3, 0xEE, 0xFE)
PINK       = RGBColor(0xEC, 0x48, 0x99)   # pink-500، لون التمييز في الموقع
PINK_PALE  = RGBColor(0xFC, 0xE7, 0xF3)
NAVY       = RGBColor(0x25, 0x45, 0x94)
GREY       = RGBColor(0x6E, 0x65, 0x79)
GREY_PALE  = RGBColor(0xF6, 0xF4, 0xFA)
GREEN      = RGBColor(0x2F, 0x7A, 0x55)
RED        = RGBColor(0xB4, 0x45, 0x3C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"
W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ────────────────────────────── أدوات ──────────────────────────────
def _rtl(par):
    """يضبط اتجاه الفقرة من اليمين لليسار، PowerPoint لا يرثه تلقائيًّا."""
    par._p.get_or_add_pPr().set('rtl', '1')


RLM = '‏'  # علامة اليمين-لليسار
# الشرطة داخل المدى (مثل 2–4) جزء من المقطع اللاتيني، وإلا انقلب الرقمان حولها
_LTR_RUN = re.compile(r'[A-Za-z0-9][A-Za-z0-9.,%/–—-]*[A-Za-z0-9%]|[A-Za-z0-9]')
_ARABIC = re.compile(r'[؀-ۿ]')


def _anchor(line):
    """يثبّت اتجاه ما بعد كل مقطع لاتيني (أرقام أو اختصارات).

    بدونها تقفز الرموز المحايدة (النقطة، %، —، ·) إلى الطرف الخاطئ
    لأن خوارزمية bidi تُلحقها بالمقطع اللاتيني بدل النصّ العربي.
    """
    return _LTR_RUN.sub(lambda m: m.group(0) + RLM, line)


def text(slide, s, x, y, w, h, size=16, bold=False, color=INK,
         align=PP_ALIGN.RIGHT, rtl=True, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, line in enumerate(str(s).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        # سطر بلا حرف عربي (مثل "40 · 150" أو "5 - 7") لاتينيّ بالكامل —
        # معاملته كعربي تقلب ترتيب الرقمين حول الفاصل.
        line_rtl = rtl and bool(_ARABIC.search(line))
        if line_rtl:
            _rtl(p)
            line = _anchor(line)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=None, line=None, radius=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def slide_light(title=None, eyebrow=None, sub=None):
    """شريحة فاتحة بترويسة موحّدة."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=PAPER)
    # شريط ذهبي رفيع أعلى اليمين، علامة الهوية
    rect(s, W - 1.55, 0, 1.55, 0.10, fill=PINK)
    if eyebrow:
        text(s, eyebrow, 0.9, 0.55, W - 1.8, 0.3, size=12, bold=True, color=PINK)
    if title:
        text(s, title, 0.9, 0.92, W - 1.8, 0.75, size=32, bold=True, color=PURPLE)
    if sub:
        text(s, sub, 0.9, 1.75, W - 1.8, 0.5, size=14, color=GREY)
    return s


def slide_dark(eyebrow=None, title=None, sub=None):
    """شريحة أرجوانية كاملة، للحظات الكبرى فقط."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=PURPLE_DEEP)
    rect(s, 0, 0, W, 0.12, fill=PINK)
    if eyebrow:
        text(s, eyebrow, 0.9, 1.5, W - 1.8, 0.3, size=13, bold=True, color=PINK_PALE)
    if title:
        text(s, title, 0.9, 1.95, W - 1.8, 1.6, size=40, bold=True, color=WHITE)
    if sub:
        text(s, sub, 0.9, 3.6, W - 1.8, 1.0, size=16, color=RGBColor(0xDC, 0xCD, 0xEC))
    return s


def card(slide, x, y, w, h, title, body, accent=PURPLE, icon=None, title_size=15):
    rect(slide, x, y, w, h, fill=PAPER, line=RGBColor(0xE3, 0xDE, 0xEB), radius=0.06)
    rect(slide, x + w - 0.06, y + 0.28, 0.06, h - 0.56, fill=accent)
    top = y + 0.28
    if icon:
        text(slide, icon, x + 0.25, top, w - 0.75, 0.4, size=20, align=PP_ALIGN.RIGHT)
        top += 0.5
    text(slide, title, x + 0.25, top, w - 0.55, 0.4, size=title_size, bold=True, color=INK)
    text(slide, body, x + 0.25, top + 0.42, w - 0.55, h - (top - y) - 0.6,
         size=11.5, color=GREY, spacing=1.3)


def kpi(slide, x, y, w, value, label, color=PURPLE, vsize=34):
    text(slide, value, x, y, w, 0.65, size=vsize, bold=True, color=color, align=PP_ALIGN.CENTER)
    text(slide, label, x, y + 0.68, w, 0.4, size=12, color=GREY, align=PP_ALIGN.CENTER)


def table(slide, x, y, w, headers, rows, col_w=None, row_h=0.42,
          head_fill=PURPLE, foot=None, size=12):
    """جدول مرسوم يدويًّا، أدقّ في التحكّم بالاتجاه والمحاذاة من جدول pptx."""
    n = len(headers)
    col_w = col_w or [w / n] * n
    # الترويسة
    rect(slide, x, y, w, 0.44, fill=head_fill)
    cx = x + w
    for i, hcell in enumerate(headers):
        cx -= col_w[i]
        align = PP_ALIGN.RIGHT if i == 0 else PP_ALIGN.CENTER
        text(slide, hcell, cx + 0.12, y + 0.11, col_w[i] - 0.24, 0.3,
             size=size - 0.5, bold=True, color=WHITE, align=align)

    yy = y + 0.44
    for ri, row in enumerate(rows):
        if ri % 2 == 1:
            rect(slide, x, yy, w, row_h, fill=GREY_PALE)
        cx = x + w
        for i, cell in enumerate(row):
            cx -= col_w[i]
            align = PP_ALIGN.RIGHT if i == 0 else PP_ALIGN.CENTER
            bold = i == 0
            text(slide, cell, cx + 0.12, yy + (row_h - 0.24) / 2, col_w[i] - 0.24, 0.26,
                 size=size, bold=bold, color=INK if i == 0 else GREY, align=align)
        yy += row_h

    if foot:
        rect(slide, x, yy, w, row_h + 0.06, fill=PURPLE_PALE)
        cx = x + w
        for i, cell in enumerate(foot):
            cx -= col_w[i]
            align = PP_ALIGN.RIGHT if i == 0 else PP_ALIGN.CENTER
            # القيم السالبة بالأحمر، أوضح من الأقواس وأسلم في النصّ العربي
            c = RED if str(cell).startswith('−') else PURPLE
            text(slide, cell, cx + 0.12, yy + 0.09, col_w[i] - 0.24, 0.3,
                 size=size + 0.5, bold=True, color=c, align=align)
        yy += row_h + 0.06
    return yy


def bars(slide, x, y, w, h, items, maxv=None, fmt=lambda v: f"{v:,.0f}",
         color=PURPLE, highlight=None, hcolor=PINK):
    """أعمدة رأسية بسيطة. items = [(label, value)]."""
    maxv = maxv or max(v for _, v in items) or 1
    n = len(items)
    gap = 0.5
    bw = (w - gap * (n - 1)) / n
    cx = x + w
    for i, (label, v) in enumerate(items):
        cx -= bw
        bh = max(0.06, (abs(v) / maxv) * h)
        c = hcolor if (highlight is not None and i == highlight) else color
        rect(slide, cx, y + h - bh, bw, bh, fill=c)
        text(slide, fmt(v), cx, y + h - bh - 0.42, bw, 0.35,
             size=15, bold=True, color=c, align=PP_ALIGN.CENTER)
        text(slide, label, cx, y + h + 0.14, bw, 0.35,
             size=12, color=GREY, align=PP_ALIGN.CENTER)
        cx -= gap


def footer(slide, num):
    text(slide, "أمانة · دراسة جدوى", 0.9, H - 0.55, 3.0, 0.3, size=9.5, color=GREY)
    text(slide, str(num), W - 1.4, H - 0.55, 0.5, 0.3, size=9.5, color=GREY,
         align=PP_ALIGN.LEFT, rtl=False)


# ═══════════════════════════ 1 · الغلاف ═══════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=RGBColor(0xFA, 0xF8, 0xFF))
rect(s, 0, 0, W, 0.16, fill=PURPLE_MID)
rect(s, W - 3.4, 0, 3.4, 0.16, fill=PINK)
rect(s, 0, H - 2.5, W, 2.5, fill=PURPLE_DEEP)

if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(W - 2.6), Inches(0.85), height=Inches(1.5))

# شارة علوية، نفس شارة بطل الموقع
rect(s, W - 6.15, 1.05, 3.3, 0.42, fill=PURPLE_PALE, radius=0.3)
text(s, "منصّة تنقّل نسائي مدعومة بالذكاء الاصطناعي", W - 6.05, 1.16, 3.1, 0.3,
     size=10.5, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

text(s, "أمانة", 0.9, 1.75, W - 3.9, 1.25, size=62, bold=True, color=PURPLE_DEEP)
text(s, "رحلتكِ… بأمان", 0.9, 3.0, W - 3.9, 0.7, size=38, bold=True, color=PURPLE_MID)
text(s, "وراحة تامّة.", 0.9, 3.72, W - 3.9, 0.7, size=38, bold=True, color=PINK)

text(s, "دراسة جدوى متكاملة · مشروع ريادي مبتكر", 0.9, H - 2.15, 7.0, 0.4, size=17,
     bold=True, color=PINK_PALE)
text(s, "مقدَّمة إلى", 0.9, H - 1.55, 5.0, 0.3, size=11, bold=True, color=PINK)
text(s, "جامعة القصيم، مركز الابتكار والملكية الفكرية", 0.9, H - 1.2, 6.5, 0.4,
     size=16, bold=True, color=WHITE)
text(s, "م. مدني كمال مدني حميدة   ·   م. عصمت عمر إبراهيم", 0.9, H - 0.72, 6.5, 0.35,
     size=12, color=RGBColor(0xD5, 0xC3, 0xF0))
text(s, "21 يوليو 2026", W - 3.3, H - 1.2, 2.4, 0.4, size=14, bold=True, color=PINK_PALE,
     align=PP_ALIGN.LEFT)

# ═══════════════════════ 2 · الفرصة والمشكلة ═══════════════════════
s = slide_light("الفرصة", "سوق قائم بلا حلّ مخصَّص",
                "المرأة في المملكة تتنقّل يوميًّا عبر تطبيقات لم تُصمَّم لها، فتدفع ثمن ذلك راحةً وأمانًا وخصوصية.")
card(s, 8.9, 2.5, 3.55, 2.05, "الخصوصية", "غالبية السائقين رجال، والرحلة اليومية تتحوّل إلى مصدر توتّر متكرّر لا خيار فيه.", PURPLE, "🔒")
card(s, 5.05, 2.5, 3.55, 2.05, "الأمان", "لا أدوات سلامة مصمَّمة للمرأة: لا مرافقة أطفال، ولا مجموعات مغلقة، ولا تحقّق مزدوج.", PURPLE, "🛡️")
card(s, 1.2, 2.5, 3.55, 2.05, "التمكين", "آلاف السعوديات يبحثن عن دخل مرن؛ قيادة المرأة فتحت الباب ولم يفتح أحد السوق.", PURPLE, "👩")
rect(s, 1.2, 4.95, 11.25, 1.05, fill=PINK_PALE, radius=0.08)
text(s, "الفجوة: لا يوجد في السوق السعودي تطبيق نقل نسائي بالكامل، من الراكبة إلى السائقة إلى فريق الدعم.",
     1.5, 5.28, 10.65, 0.45, size=15, bold=True, color=PINK, align=PP_ALIGN.CENTER)
footer(s, 2)

# ═══════════════════════ 3 · حجم السوق ═══════════════════════
s = slide_light("حجم السوق", "سوق ينمو 13% سنويًّا، ونصفه بلا خدمة",
                "قطاع طلب سيارات الأجرة في المملكة، وأرقام الطلب النسائي التي تثبت جاهزية السوق.")
for i, (v, l, c) in enumerate([("1.9 مليار $", "حجم السوق 2025", PURPLE),
                               ("5.6 مليار $", "المتوقَّع 2034", PURPLE),
                               ("13.03%", "نموّ سنوي مركّب", PINK),
                               ("54.26%", "نموّ الرحلات سنويًّا", PINK)]):
    x = W - 1.2 - (i + 1) * 2.83 + 0.15
    rect(s, x, 2.5, 2.62, 1.5, fill=PAPER, line=RGBColor(0xE3, 0xDE, 0xEB), radius=0.08)
    rect(s, x, 2.5, 2.62, 0.07, fill=c)
    text(s, v, x, 2.82, 2.62, 0.5, size=22, bold=True, color=c, align=PP_ALIGN.CENTER)
    text(s, l, x, 3.42, 2.62, 0.35, size=12, color=GREY, align=PP_ALIGN.CENTER)

text(s, "سوقا الانطلاق: بريدة أوّلًا، ثمّ الرياض", 1.2, 4.25, 11.25, 0.35,
     size=15, bold=True, color=PURPLE)
table(s, 1.2, 4.7, 11.25,
      ["المؤشّر", "بريدة — الانطلاق", "الرياض — التوسّع", "لماذا هذا الترتيب"],
      [["عدد السكّان", "~700 ألف", "~8 ملايين", "نبدأ صغيرًا لنتعلّم، ثمّ ندخل الأكبر بثقة"],
       ["نسبة النساء", "~45%", "~42%", "قاعدة الطلب في المدينتين"],
       ["الموظفات والطالبات", "~250 ألف", "~1.4 مليون", "الشريحة الأعلى تكرارًا للرحلات"],
       ["ميزة إضافية", "مقرّ جامعة القصيم", "أكبر سوق بالمملكة", "الحاضنة تجاورنا في سوق الاختبار"]],
      col_w=[2.35, 2.15, 2.15, 4.6], row_h=0.44, size=11.5)
footer(s, 3)

# ═══════════════════════════ 4 · الحلّ ═══════════════════════════
s = slide_dark("الحلّ", "منظومة تنقّل نسائية\nمدعومة بالذكاء الاصطناعي",
               "ليست تطبيق نقل فحسب، بل بيئة كاملة صُمّمت حول أمان المرأة وخصوصيتها وتمكينها اقتصاديًّا.")
for i, (v, l) in enumerate([("3", "تطبيقات متكاملة"), ("100%", "نسائية بالكامل"),
                            ("AI", "ذكاء اصطناعي مدمج"), ("جاهز", "نموذج أوّلي يعمل")]):
    x = W - 1.0 - (i + 1) * 2.85 + 0.15
    text(s, v, x, 5.05, 2.6, 0.6, size=30, bold=True, color=PINK_PALE, align=PP_ALIGN.CENTER)
    text(s, l, x, 5.72, 2.6, 0.35, size=12.5, color=RGBColor(0xD2, 0xC0, 0xE6), align=PP_ALIGN.CENTER)
footer(s, 4)

# ══════════════════ 5 · مصفوفة التفرّد ══════════════════
s = slide_light("الميزة التنافسية", "ثماني مزايا لا يملكها أيّ منافس مجتمعة",
                "كل ميزة تحلّ مشكلة حقيقية، وأربع منها غير موجودة في السوق السعودي إطلاقًا.")
feats = [
    ("🧠", "مساعد السفر الذكي", "يقترح وجهات بناءً على اهتمامات الراكبة وتاريخ رحلاتها."),
    ("📈", "التسعير الديناميكي", "سعر شفّاف يعكس الطلب لحظيًّا، ويُعرض قبل التأكيد."),
    ("🌿", "البصمة الكربونية", "حساب التوفير في الانبعاثات ومكافآت خضراء للراكبات."),
    ("🛡️", "التقييم المزدوج", "تقييم الراكبة والسائقة معًا، حماية متبادلة ترفع الجودة."),
    ("👩‍👧", "مرافقة الأطفال", "سائقات مدرَّبات على تركيب مقاعد الأطفال، لا يقدّمها منافس."),
    ("🔒", "المجموعات المغلقة", "رحلات تشاركية خاصّة بين زميلات العمل أو الدراسة."),
    ("👥", "النقل التشاركي", "مشاركة التكلفة بين الراكبات، خفض يصل إلى 60%."),
    ("🗺️", "الرحلات السياحية", "مسارات معدّة لاستكشاف معالم المدن، شريحة الزائرات."),
]
for i, (ic, t, b) in enumerate(feats):
    col, row = i % 4, i // 4
    x = W - 1.2 - (col + 1) * 2.83 + 0.15
    y = 2.5 + row * 2.05
    card(s, x, y, 2.62, 1.9, t, b, PURPLE if row == 0 else PINK, ic, title_size=13)
footer(s, 5)

# ══════════════════ 6 · تحليل المنافسين ══════════════════
s = slide_light("المنافسة", "لا أحد يجمع التخصّص والابتكار معًا",
                "المنافس النسائي الحالي يفتقر للابتكار، والمنافس المبتكر ليس نسائيًّا بالكامل.")
table(s, 1.2, 2.5, 11.25,
      ["المنافس", "الحصّة", "نقطة القوّة", "نقطة الضعف", "التهديد"],
      [["وصليني", "~15% نسائيًّا", "أوّل تطبيق نسائي، 60 مدينة", "بلا ذكاء اصطناعي ولا استدامة", "متوسّط"],
       ["مريني", "~10% نسائيًّا", "تركيز على الخصوصية", "مرحلة تطوير، تركيز ضيّق", "منخفض"],
       ["أوبر «سائقات»", "~5% من أوبر", "علامة عالمية وشبكة ضخمة", "خدمة عامّة لا منصّة نسائية", "مرتفع"],
       ["كريم · جيني وغيرهما", "~70%", "تغطية وأسعار تنافسية", "لا خدمة نسائية حصرية", "منخفض"]],
      col_w=[2.65, 1.75, 2.95, 2.95, 0.95], row_h=0.58, size=12)
rect(s, 1.2, 5.55, 11.25, 1.0, fill=PINK_PALE, radius=0.08)
text(s, "استراتيجية التفوّق: مزايا حصرية يصعب تقليدها: مرافقة الأطفال، البصمة الكربونية، المجموعات المغلقة\nمع توسّع سريع في المدن الرئيسية قبل أن يتحرّك اللاعبون الكبار.",
     1.5, 5.75, 10.65, 0.65, size=13, bold=True, color=PINK, align=PP_ALIGN.CENTER)
footer(s, 6)

# ═══════════════════ 7 · فرضيات النموّ ═══════════════════
s = slide_light("أساس الحساب", "فرضيات النموّ",
                "كل رقم في هذا العرض مشتقّ من هذه الفرضيات، تغييرها يغيّر النتائج بالتناسب.")
table(s, 1.2, 2.55, 11.25,
      ["المؤشّر", "السنة 1، الإطلاق", "السنة 2، النموّ", "السنة 3، التوسّع"],
      [["المدن المغطّاة", "بريدة ثمّ الرياض", "3", "5 - 7"],
       ["الرحلات سنويًّا", "30,000", "300,000", "900,000"],
       ["مستخدمات نشطات شهريًّا", "1,200", "10,000", "30,000"],
       ["سائقات نشطة · مسجَّلة", "40 · 150", "420 · 1,200", "1,250 · 4,000"],
       ["متوسّط قيمة الرحلة", "30 ريال", "30 ريال", "30 ريال"]],
      col_w=[3.75, 2.5, 2.5, 2.5], row_h=0.52, size=13)
rect(s, 1.2, 5.6, 11.25, 0.95, fill=PURPLE_PALE, radius=0.1)
text(s, "الأرقام مترابطة حسابيًّا: 2.5 رحلة شهريًّا لكل مستخدمة نشطة، و60 رحلة شهريًّا لكل سائقة نشطة.\nنبدأ من بريدة في الشهر السابع بدعم حاضنة جامعة القصيم، وندخل الرياض في الشهر العاشر من السنة نفسها.",
     1.5, 5.78, 10.65, 0.65, size=12.5, color=PURPLE, align=PP_ALIGN.CENTER)
footer(s, 7)

# ══════════════ 8 · الأطروحة: تكلفة الرحلة ══════════════
s = slide_light("الأطروحة الاستثمارية", "التكلفة التقنية لكل رحلة تنخفض 92%",
                "بنية سحابية مُدارة وفريق يكتمل تدريجيًّا: التكلفة تكاد تثبت بينما الرحلات تتضاعف ثلاثين مرّة.")
bars(s, 2.6, 2.75, 8.4, 2.45,
     [("السنة 1\n30 ألف رحلة", 24), ("السنة 2\n300 ألف رحلة", 5), ("السنة 3\n900 ألف رحلة", 2)],
     maxv=24, fmt=lambda v: f"{v} ر.س", highlight=2)
rect(s, 3.6, 6.05, 6.1, 0.62, fill=PINK_PALE, radius=0.2)
text(s, "▼  24 ← 2 ريال   ·   الميزانية التقنية ÷ عدد الرحلات، بلا تقريب", 3.6, 6.24, 6.1, 0.35,
     size=13.5, bold=True, color=PINK, align=PP_ALIGN.CENTER)
footer(s, 8)

# ══════════════════ 9 · الميزانية التقنية ══════════════════
s = slide_light("الاستثمار التقني", "الميزانية التقنية، ثلاث سنوات",
                "خدمات سحابية مُدارة بالكامل: لا خوادم مملوكة ولا مراكز بيانات ولا فريق تشغيل.")
table(s, 1.2, 2.5, 11.25,
      ["البند", "السنة 1", "السنة 2", "السنة 3", "الإجمالي"],
      [["البنية التحتية", "15,050", "63,000", "178,000", "256,050"],
       ["الذكاء الاصطناعي", "2,250", "116,000", "247,000", "365,250"],
       ["الفريق التقني", "530,600", "1,141,600", "1,220,500", "2,892,700"],
       ["الأمن والأدوات", "100,000", "115,000", "175,000", "390,000"],
       ["احتياطي طوارئ 10%", "64,790", "143,560", "182,050", "390,400"]],
      col_w=[3.05, 2.05, 2.05, 2.05, 2.05], row_h=0.5, size=12.5,
      foot=["الإجمالي", "712,690", "1,579,160", "2,002,550", "4,294,400"])
rect(s, 1.2, 6.0, 11.25, 0.75, fill=PINK_PALE, radius=0.1)
text(s, "الفريق يكتمل تدريجيًّا: السنة الأولى بنصف الطاقم، لا نوظّف ثمانية أفراد قبل أن توجد رحلات يخدمونها.",
     1.5, 6.22, 10.65, 0.4, size=12.5, color=PINK, align=PP_ALIGN.CENTER)
footer(s, 9)

# ═══════════════════ 10 · الفريق التقني ═══════════════════
s = slide_light("الهيكل", "فريق يكتمل مع الطلب، لا قبله",
                "المؤسّسان هما نواة الفريق التقني: قيادة تقنية وتشغيل بلا تكلفة توظيف خارجية.")
table(s, 1.2, 2.55, 11.25,
      ["الدور", "متى ينضمّ", "العدد", "السنة 1", "السنة 2", "السنة 3"],
      [["قائد تقني: م. مدني، مؤسّس", "من الشهر 1", "1", "120,000", "277,300", "296,400"],
       ["مطوّر Full-stack", "من الشهر 1", "1", "187,200", "200,400", "214,200"],
       ["مطوّر مساعد", "الشهر 4 ← 2 بالسنة 2", "1 ← 2", "86,400", "246,600", "263,500"],
       ["مطوّر واجهات UI/UX", "عقد 6 أشهر ← دوام", "1", "65,000", "140,000", "150,000"],
       ["دعم فنّي", "الشهر 6 ← 3 بالسنة 2", "2 ← 3", "72,000", "277,300", "296,400"]],
      col_w=[3.35, 2.45, 1.2, 1.4, 1.42, 1.43], row_h=0.5, size=11.5,
      foot=["الإجمالي", "", "6 ← 8", "530,600", "1,141,600", "1,220,500"])
rect(s, 1.2, 6.15, 11.25, 0.7, fill=PURPLE_PALE, radius=0.1)
text(s, "راتب المؤسّس في السنة الأولى مخفَّض إلى 10,000 ريال شهريًّا، التزام شخصي برأس المال قبل أن نطلبه من غيرنا.",
     1.5, 6.35, 10.65, 0.35, size=12, color=PURPLE, align=PP_ALIGN.CENTER)
footer(s, 10)

# ══════════════════ 11 · التوقّعات المالية ══════════════════
s = slide_light("النموذج المالي", "التعادل في السنة الثانية، والربحية في الثالثة",
                "العمولة هي العمود الفقري للإيراد. الأرقام بملايين الريالات، وكل صفّ مشتقّ من فرضيات شريحة أساس الحساب.")
table(s, 1.2, 2.45, 11.25,
      ["المؤشّر", "السنة 1", "السنة 2", "السنة 3"],
      [["الرحلات سنويًّا", "30,000", "300,000", "900,000"],
       ["إجمالي قيمة الرحلات GMV", "0.90", "9.00", "27.00"],
       ["الإيرادات: عمولة 20% وباقات وخدمات", "0.25", "2.25", "6.75"],
       ["التكاليف الكلّية: تقنية وتشغيلية", "1.11", "2.48", "4.40"]],
      col_w=[4.65, 2.2, 2.2, 2.2], row_h=0.46, size=13,
      foot=["صافي الربح أو الخسارة", "−0.86", "−0.23", "2.35"])
for i, (v, l, c) in enumerate([("السنة 3", "ربحية كاملة بـ2.35 مليون", GREEN),
                               ("السنة 2", "التعادل في الربع الأخير", PINK),
                               ("السنة 1", "بناء السوق والشبكة", PURPLE)]):
    x = 1.2 + i * 3.79
    rect(s, x, 5.45, 3.65, 1.0, fill=PAPER, line=RGBColor(0xE3, 0xDE, 0xEB), radius=0.08)
    rect(s, x, 5.45, 3.65, 0.07, fill=c)
    text(s, v, x, 5.66, 3.65, 0.32, size=15, bold=True, color=c, align=PP_ALIGN.CENTER)
    text(s, l, x, 6.02, 3.65, 0.3, size=11.5, color=GREY, align=PP_ALIGN.CENTER)
text(s, "الخسارة في أوّل سنتين مُقدَّرة ومموَّلة سلفًا، وهي 7% فقط من إجمالي قيمة الرحلات التراكمية.",
     1.2, 6.6, 11.25, 0.35, size=12, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 11)

# ═══════════ 12 · اقتصاديات الرحلة الواحدة ═══════════
s = slide_dark("اقتصاديات الوحدة", "2.60 ريال هامشًا صافيًا\nمن كل رحلة في السنة الثالثة")
for i, (v, l, c) in enumerate([("7.50", "إيراد المنصّة من الرحلة", PINK_PALE),
                               ("4.90", "التكلفة الكلّية للرحلة", RGBColor(0xC9, 0xB6, 0xDE)),
                               ("2.60", "صافي الهامش", WHITE)]):
    x = W - 1.0 - (i + 1) * 3.75 + 0.15
    text(s, v, x, 4.3, 3.5, 0.8, size=42, bold=True, color=c, align=PP_ALIGN.CENTER)
    text(s, l + " · ريال", x, 5.2, 3.5, 0.35, size=13, color=RGBColor(0xD2, 0xC0, 0xE6),
         align=PP_ALIGN.CENTER)
text(s, "هامش صافٍ 35%، بعد احتساب كامل تكلفة الفريق والبنية والتشغيل والتسويق داخل الرقم، لا قبلها.",
     0.9, 6.1, W - 1.8, 0.4, size=13, color=PINK_PALE, align=PP_ALIGN.CENTER)
footer(s, 12)

# ═══════════ 13 · باقة التنقّل ═══════════
s = slide_light("منتج الولاء", "«باقة التنقّل»، 149 ريال شهريًّا",
                "أُعيد تصميم الباقة لتكون رابحة بذاتها: خصم على الرحلات، لا رحلات مجّانية.")
table(s, 1.2, 2.6, 7.4, ["البند، لمشتركة تنفّذ 16 رحلة شهريًّا", "المبلغ"],
      [["قيمة الرحلات بالسعر المعتاد", "480 ريال"],
       ["ما تدفعه المشتركة بعد خصم 25%", "360 ريال"],
       ["رسوم الاشتراك الشهري", "149 ريال"],
       ["إجمالي ما يصل المنصّة", "509 ريال"],
       ["حصّة السائقة، 80% من السعر المعتاد", "−384 ريال"]],
      col_w=[5.4, 2.0], row_h=0.5, size=12.5,
      foot=["صافي المنصّة شهريًّا", "125 ريال"])
card(s, 8.85, 2.6, 3.6, 1.75, "لماذا هي رابحة",
     "125 ريالًا شهريًّا مقابل 96 ريالًا فقط من راكبة غير مشتركة تنفّذ العدد نفسه، الباقة تزيد الهامش ولا تأكله.", GREEN, "✅")
card(s, 8.85, 4.55, 3.6, 1.75, "دورها الحقيقي",
     "أداة ولاء وتدفّق نقدي مقدَّم، لا عمود إيراد: نُقدِّرها بـ4% من المستخدمات النشطات فقط.", PURPLE, "🔁")
rect(s, 1.2, 6.25, 7.4, 0.55, fill=PINK_PALE, radius=0.1)
text(s, "التصميم السابق، 10 رحلات مجّانية، كان يخسر 91 ريالًا شهريًّا عن كل مشتركة.",
     1.4, 6.4, 7.0, 0.35, size=12, bold=True, color=PINK, align=PP_ALIGN.CENTER)
footer(s, 13)

# ═══════════════════ 14 · خارطة الطريق ═══════════════════
s = slide_light("التنفيذ", "خارطة الطريق، مدينتان في السنة الأولى",
                "نختبر النموذج في بريدة حيث تجاورنا الحاضنة، ثمّ ندخل الرياض بمنتج مُثبَت لا بفرضية.")
phases = [
    ("1–4", "التأسيس والترخيص الريادي", "خطاب تزكية الجامعة، خطاب عدم الممانعة، رخصة «ريادي»، السجل التجاري", "25,000"),
    ("2–6", "التطوير التقني MVP", "تطبيقات Android و iOS، الربط مع منصّة هيئة النقل، واجهات احترافية", "150,000"),
    ("5–6", "استقطاب السائقات، بريدة", "أوّل 100 سائقة وتأهيلهنّ على القيادة الآمنة وخدمة العملاء", "60,000"),
    ("7", "الإطلاق التجريبي، بريدة", "التشغيل مع 50 سائقة وجمع الملاحظات الميدانية وتصحيح المسار", "25,000"),
    ("8–9", "الإطلاق الرسمي، بريدة", "حملة موجَّهة بالتعاون مع جامعة القصيم وإطلاق «باقة التنقّل»", "70,000"),
    ("10–12", "دخول الرياض", "استقطاب 100 سائقة وحملة تسويق في أكبر سوق بالمملكة", "190,000"),
]
yy = 2.5
for i, (m, t, d, c) in enumerate(phases):
    if i % 2 == 1:
        rect(s, 1.2, yy, 11.25, 0.62, fill=GREY_PALE)
    rect(s, 11.65, yy + 0.11, 0.72, 0.4, fill=PURPLE_PALE, radius=0.2)
    text(s, m, 11.65, yy + 0.2, 0.72, 0.25, size=11.5, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    text(s, t, 8.6, yy + 0.09, 2.9, 0.28, size=13.5, bold=True, color=INK)
    text(s, d, 3.2, yy + 0.13, 5.25, 0.3, size=11.5, color=GREY)
    text(s, c + " ريال", 1.3, yy + 0.12, 1.75, 0.3, size=12.5, bold=True, color=PINK, align=PP_ALIGN.LEFT)
    yy += 0.62
rect(s, 1.2, yy + 0.12, 11.25, 0.6, fill=PURPLE)
text(s, "تكلفة إطلاق المدينتين، عدا الضمان المسترَدّ", 7.4, yy + 0.29, 4.8, 0.3,
     size=13.5, bold=True, color=WHITE)
text(s, "520,000 ريال", 1.3, yy + 0.29, 3.0, 0.3, size=15, bold=True, color=PINK_PALE, align=PP_ALIGN.LEFT)
footer(s, 14)

# ════════════════ 15 · رأس المال والتمويل ════════════════
s = slide_light("الطلب", "1.5 مليون ريال، تمويل حتى التعادل",
                "إطلاق المدينتين يحتاج 770 ألف ريال؛ والوصول إلى التعادل يحتاج تغطية عجز السنتين الأولى والثانية.")
table(s, 1.2, 2.6, 7.4, ["أوجه الاستخدام", "المبلغ", "النسبة"],
      [["الضمان المالي لهيئة النقل، مسترَدّ", "250,000", "17%"],
       ["تغطية عجز التشغيل حتى التعادل", "545,000", "36%"],
       ["تطوير المنتج والبنية التقنية", "250,000", "17%"],
       ["تسويق بريدة والرياض", "180,000", "12%"],
       ["استقطاب السائقات وتأهيلهنّ", "140,000", "9%"],
       ["التأسيس والتراخيص واحتياطي", "135,000", "9%"]],
      col_w=[4.0, 2.0, 1.4], row_h=0.45, size=12,
      foot=["الإجمالي", "1,500,000", "100%"])
card(s, 8.85, 2.6, 3.6, 1.75, "مصادر الدعم المستهدفة",
     "كفالة · رواد التقنية · بنك التنمية الاجتماعية · حاضنة ابتكار بجامعة القصيم", PINK, "🤝")
card(s, 8.85, 4.55, 3.6, 1.75, "العائد المتوقّع",
     "استرداد كامل خلال السنة الثالثة، وصافي ربح 2.35 مليون ريال في السنة الثالثة وحدها.", GREEN, "📈")
rect(s, 1.2, 6.25, 7.4, 0.55, fill=PURPLE_PALE, radius=0.1)
text(s, "الاستنزاف النقدي الفعلي 1.25 مليون، إذ يُستردّ الضمان البالغ 250 ألفًا.",
     1.4, 6.4, 7.0, 0.35, size=12, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
footer(s, 15)

# ════════════════ 16 · المسار النظامي ════════════════
s = slide_light("المسار النظامي", "الترخيص الريادي، مسار رسمي مكتمل",
                "رخصة تصدرها وزارة الاستثمار لروّاد الأعمال غير السعوديين: ملكية 100% وإقامة مستقلّة.")
steps = [("1", "خطاب التزكية", "من مركز الابتكار بجامعة القصيم، الشرط الأهمّ لقبول الطلب", "2–4 أسابيع"),
         ("2", "خطاب عدم الممانعة", "من الكفيل الحالي لتحويل الإقامة", "2–5 أيام"),
         ("3", "رخصة «ريادي»", "التقديم عبر بوابة وزارة الاستثمار MISA", "5–10 أيام عمل"),
         ("4", "تحويل الإقامة", "الإقامة الريادية عبر الجوازات، تشمل العائلة", "2–4 أسابيع"),
         ("5", "تأسيس الشركة", "السجل التجاري، العلامة التجارية، ترخيص هيئة النقل", "3–7 أيام")]
yy = 2.5
for i, (n, t, d, dur) in enumerate(steps):
    if i % 2 == 1:
        rect(s, 1.2, yy, 11.25, 0.66, fill=GREY_PALE)
    rect(s, 11.72, yy + 0.13, 0.6, 0.4, fill=PURPLE_PALE, radius=0.3)
    text(s, n, 11.72, yy + 0.22, 0.6, 0.25, size=12, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    text(s, t, 8.5, yy + 0.11, 3.1, 0.3, size=13.5, bold=True, color=INK)
    text(s, d, 3.0, yy + 0.15, 5.35, 0.3, size=11.5, color=GREY)
    text(s, dur, 1.3, yy + 0.14, 1.6, 0.3, size=12, bold=True, color=PINK, align=PP_ALIGN.LEFT)
    yy += 0.6

# الرسوم الرسمية الفعلية، لا تقديرات
rect(s, 1.2, 5.52, 11.25, 0.52, fill=PURPLE_PALE, radius=0.1)
text(s, "الرسوم الرسمية للتأسيس: رخصة ريادي 2,000 · سجل تجاري 1,200 · توثيق ونشر عقد التأسيس 500 · ضريبة قيمة مضافة 15% · غرفة تجارية 300 فأعلى",
     1.4, 5.65, 10.85, 0.35, size=11.5, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

for i, (t, b) in enumerate([("ملكية 100%", "بلا شريك محلّي"),
                            ("2,000 ريال/سنة", "رسوم أوّل 3 سنوات"),
                            ("إعفاء من نطاقات", "3 سنوات تشغيلية"),
                            ("إقامة عائلية", "الزوج/ة والأبناء والوالدان")]):
    x = W - 1.2 - (i + 1) * 2.83 + 0.15
    rect(s, x, 6.18, 2.62, 0.7, fill=PINK_PALE, radius=0.1)
    text(s, t, x, 6.27, 2.62, 0.3, size=13, bold=True, color=PINK, align=PP_ALIGN.CENTER)
    text(s, b, x, 6.58, 2.62, 0.26, size=10, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 16)

# ═════════════ 17 · التوافق مع رؤية 2030 ═════════════
s = slide_light("التوافق الوطني", "أمانة ورؤية 2030",
                "المشروع يخدم أربعة مستهدفات وطنية في آنٍ واحد، وهذا ما يجعله مؤهّلًا للدعم.")
v30 = [("👩‍💼", "تمكين المرأة", "4,000 سائقة مؤهَّلة ومسجَّلة، منهنّ 1,250 بدخل نشط بحلول السنة الثالثة."),
       ("💡", "الابتكار التقني", "ذكاء اصطناعي وتسعير ديناميكي مبنيّ محلّيًّا بكوادر وطنية."),
       ("🌿", "الاستدامة", "نقل تشاركي يخفض ~2.5 كجم كربون لكل رحلة ويقيسها ويكافئ عليها."),
       ("🏙️", "تنمية المناطق", "الإطلاق من بريدة يضع القصيم في مقدّمة الخدمة، لا في آخر قائمة التوسّع.")]
for i, (ic, t, b) in enumerate(v30):
    col, row = i % 2, i // 2
    x = W - 1.2 - (col + 1) * 5.65 + 0.25
    y = 2.55 + row * 2.05
    card(s, x, y, 5.4, 1.85, t, b, PURPLE if row == 0 else PINK, ic, title_size=15)
footer(s, 17)

# ══════════════════ 18 · إدارة المخاطر ══════════════════
s = slide_light("الحوكمة", "إدارة المخاطر، لكل احتمال خطّة بديلة",
                "لا نعرض المخاطر فحسب، بل نعرض ما سنفعله فعليًّا إن تحقّقت.")
table(s, 1.2, 2.5, 11.25, ["المخاطرة", "استراتيجية التخفيف", "خطّة الطوارئ"],
      [["سوقية، تباطؤ النموّ في إحدى المدينتين",
        "مدينتان مختلفتان في الحجم والسلوك، لا رهان على سوق واحد",
        "تحويل الميزانية والسائقات إلى المدينة الأسرع استجابة"],
       ["قانونية، رفض التزكية أو الترخيص",
        "تواصل مبكر مع الجهات، ومستشار قانوني مدرج في الميزانية",
        "حاضنات بديلة، واعد أو مسك، أو شريك سعودي، أو تأشيرة استثمارية"],
       ["مالية، تجاوز التكاليف أو نقص الإيراد",
        "احتياطي 10%، وتوظيف تدريجي يجعل أكبر بند تكلفة قابلًا للتأجيل",
        "تأجيل التوظيف وخفض المصاريف والتركيز على النموّ العضوي"],
       ["تشغيلية، نقص السائقات أو ضعف الجودة",
        "حوافز ذروة، تدريب مكثّف قبل الإطلاق، تقييم فوري",
        "عروض تشجيعية للسائقات والتعاقد مع شركات توظيف"],
       ["تنافسية، دخول لاعب أكبر تمويلًا",
        "مزايا حصرية يصعب تقليدها ومجتمع نسائي مغلق",
        "تسريع التوسّع للمدن الرئيسية والتمايز بالميزات المجتمعية"],
       ["تنظيمية، ارتفاع رسوم الرخصة بعد الإعفاء",
        "الإعفاء يغطّي السنوات الثلاث الأولى، وهي كامل مدى هذه الخطة",
        "الرسوم بعدها تُحتسب من إيراد السنة الثالثة، لا من رأس المال"],
       ["تقنية، تعطّل أو اختراق",
        "بنية سحابية مُدارة، نسخ احتياطي يومي، اختبار اختراق سنوي",
        "دعم فنّي على مدار الساعة واستعادة من النسخ الاحتياطية"]],
      col_w=[3.1, 4.05, 4.1], row_h=0.56, size=10.5)
footer(s, 18)

# ═════════════════ 19 · لماذا الآن ═════════════════
s = slide_dark("الخلاصة", "لماذا أمانة، ولماذا الآن؟")
whys = [("السوق مفتوح", "لا منافس نسائي متكامل في المملكة حتى اليوم"),
        ("المنتج جاهز", "نموذج أوّلي يعمل فعليًّا، ليس فكرة على ورق"),
        ("الاقتصاد سليم", "هامش صافٍ 35% للرحلة وتعادل في السنة الثانية"),
        ("التوقيت مثالي", "رؤية 2030 وبرامج الدعم في ذروة تفعيلها")]
for i, (t, b) in enumerate(whys):
    col, row = i % 2, i // 2
    x = W - 1.0 - (col + 1) * 5.7 + 0.2
    y = 3.5 + row * 1.4
    text(s, "◆  " + t, x, y, 5.3, 0.4, size=19, bold=True, color=PINK_PALE)
    text(s, b, x + 0.35, y + 0.45, 4.95, 0.4, size=13, color=RGBColor(0xD2, 0xC0, 0xE6))
footer(s, 19)

# ═══════════════════ 20 · التواصل ═══════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
rect(s, 0, 0, W, 0.14, fill=PINK)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(W / 2 - 0.7), Inches(1.15), height=Inches(1.4))
text(s, "أمانة", 0, 2.75, W, 0.9, size=46, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
text(s, "رحلتكِ… بأمان وراحة تامّة", 0, 3.7, W, 0.5, size=19, color=GREY, align=PP_ALIGN.CENTER)
rect(s, 3.4, 4.6, 6.5, 0.04, fill=PINK)

text(s, "م. مدني كمال مدني حميدة", 6.95, 5.0, 5.0, 0.35, size=15, bold=True, color=INK)
text(s, "00966579597906", 6.95, 5.42, 5.0, 0.32, size=13, color=GREY)
text(s, "medani7@gmail.com", 6.95, 5.8, 5.0, 0.35, size=13.5, bold=True, color=PURPLE,
     align=PP_ALIGN.RIGHT, rtl=False)

text(s, "م. عصمت عمر إبراهيم", 1.4, 5.0, 5.0, 0.35, size=15, bold=True, color=INK)
text(s, "00966533362837", 1.4, 5.42, 5.0, 0.32, size=13, color=GREY)
text(s, "esmat369@hotmail.com", 1.4, 5.8, 5.0, 0.35, size=13.5, bold=True, color=PURPLE,
     align=PP_ALIGN.RIGHT, rtl=False)

prs.save(OUT)
print("saved:", OUT, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
